# -*- coding: utf-8 -*-
"""
CA-ISO 电价价差预测 · 讲解版 PPT 生成脚本
运行：python code/deliverables/make_ppt.py
输出：code/deliverables/讲解版.pptx
"""
import os

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------
# 路径
# ----------------------------------------------------------------------------
BASE = r"D:\code\pyCode\CA-电力交易预测"
CHART_DIR = os.path.join(BASE, "code", "deliverables", "charts")
OUT_PATH = os.path.join(BASE, "code", "deliverables", "讲解版.pptx")
ARB_PNG = os.path.join(BASE, "code", "data", "arb_curve.png")
os.makedirs(CHART_DIR, exist_ok=True)

# ----------------------------------------------------------------------------
# 配色（深蓝主色 + 白底）
# ----------------------------------------------------------------------------
NAVY      = RGBColor(0x1F, 0x4E, 0x79)   # 深蓝主色
NAVY_DARK = RGBColor(0x14, 0x36, 0x52)   # 封面底
BLUE_MED  = RGBColor(0x2E, 0x75, 0xB6)   # 中蓝
GOLD      = RGBColor(0xC9, 0xA2, 0x27)   # 金色点缀
RED       = RGBColor(0xC0, 0x00, 0x00)
GREEN     = RGBColor(0x53, 0x8D, 0x4E)
DARK      = RGBColor(0x33, 0x33, 0x33)
GRAY      = RGBColor(0x59, 0x59, 0x59)
LIGHT     = RGBColor(0xF2, 0xF2, 0xF2)
LIGHT_BLUE= RGBColor(0xDE, 0xEB, 0xF7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

NAVY_HEX   = "#1F4E79"
BLUE_HEX   = "#2E75B6"
GOLD_HEX   = "#C9A227"
RED_HEX    = "#C00000"
GREEN_HEX  = "#538D4E"
GRAY_HEX   = "#8C8C8C"

# 图表中文字体
plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "DengXian", "sans-serif"]
plt.rcParams["axes.unicode_minus"] = False

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.55)
CONTENT_W = Emu(int(SLIDE_W) - 2 * int(MARGIN))

# ----------------------------------------------------------------------------
# 图表 1：特征重要性（top 11 条形图）+ 类别汇总（环形图）
# ----------------------------------------------------------------------------
def make_feature_chart(path):
    feats = [
        ("da_lag2（D-2 日前价）", 15.8),
        ("spread_std14（近14日价差波动）", 13.1),
        ("spread_std7（近7日价差波动）", 7.8),
        ("node（节点）", 6.9),
        ("month_next（月份）", 5.6),
        ("da_lag1（D-1 日前价）", 5.5),
        ("peer_rtpd_lag1（联动节点实时价滞后）", 5.0),
        ("load_2da_next（D+1 负荷预测）", 5.0),
        ("da_lag7（D-7 日前价）", 4.5),
        ("rtpd_lag1（D-1 实时价）", 4.5),
        ("ssrd_next（D+1 太阳辐射）", 4.1),
        ("spread_lag7（D-7 价差）", 3.5),
    ]
    cats = [("历史滞后", 44.4), ("滚动统计（波动）", 24.1), ("未来预报", 22.9), ("日历", 9.7), ("节点", 6.9)]

    labels = [f[0] for f in feats][::-1]
    values = [f[1] for f in feats][::-1]
    colors = [GOLD_HEX if v >= 13 else NAVY_HEX for v in values]

    fig = plt.figure(figsize=(11.6, 4.75), dpi=200)
    ax1 = fig.add_axes([0.02, 0.02, 0.55, 0.94])
    ax2 = fig.add_axes([0.66, 0.02, 0.33, 0.94])

    bars = ax1.barh(labels, values, color=colors, height=0.62)
    ax1.set_xlim(0, 18)
    ax1.set_xlabel("特征重要性（占比 %）", fontsize=11)
    ax1.set_title("CatBoost 方向分类器 · Top 特征重要性", fontsize=13, fontweight="bold", color=NAVY_HEX, loc="left", pad=10)
    ax1.tick_params(axis="y", labelsize=10.5)
    ax1.tick_params(axis="x", labelsize=10)
    for bar, v in zip(bars, values):
        ax1.text(bar.get_width() + 0.6, bar.get_y() + bar.get_height() / 2,
                 f"{v:.1f}%", va="center", ha="left", fontsize=10, color="#333333")
    ax1.spines[["top", "right"]].set_visible(False)
    ax1.grid(axis="x", linestyle="--", alpha=0.35)

    cvals = [c[1] for c in cats]
    clabels = [c[0] for c in cats]
    ccolors = [NAVY_HEX, "#7F9DC0", BLUE_HEX, GOLD_HEX, GRAY_HEX]
    wedges, _ = ax2.pie(cvals, colors=ccolors, startangle=90, counterclock=False,
                        wedgeprops=dict(width=0.42, edgecolor="white", linewidth=2))
    ax2.text(0, 0.08, "类别汇总", ha="center", va="center", fontsize=12, fontweight="bold", color=NAVY_HEX)
    ax2.text(0, -0.13, "按占比", ha="center", va="center", fontsize=10, color=GRAY_HEX)
    ax2.legend(wedges, [f"{l}  {v:.1f}%" for l, v in zip(clabels, cvals)],
               loc="center left", bbox_to_anchor=(1.02, 0.5), fontsize=10.5, frameon=False)

    fig.savefig(path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print("chart saved:", path)


# ----------------------------------------------------------------------------
# 图表 2：基线对比（横向柱状图）
# ----------------------------------------------------------------------------
def make_baseline_chart(path):
    rows = [
        ("本工程（单边策略）", +2243,   NAVY_HEX, True),
        ("不交易",              0,       GRAY_HEX, False),
        ("反向策略（镜像上界）", +128259, GOLD_HEX, False),
        ("全交易（无观望/风控）", -128259, RED_HEX,  False),
    ]
    # 纵向排序：大的在上
    rows = sorted(rows, key=lambda r: r[1])
    labels = [r[0] for r in rows]
    values = [r[1] for r in rows]
    colors = [r[2] for r in rows]

    fig, ax = plt.subplots(figsize=(11.6, 4.3), dpi=200)
    bars = ax.barh(labels, values, color=colors, height=0.6)
    ax.axvline(0, color="#999999", linewidth=1)
    ax.set_xlabel("test 模拟套利收益（美元，65 天）", fontsize=11)
    ax.set_title("基线对比：test 模拟套利收益", fontsize=13, fontweight="bold", color=NAVY_HEX, loc="left", pad=10)
    ax.tick_params(axis="y", labelsize=11)
    ax.tick_params(axis="x", labelsize=10)

    xmax = max(abs(v) for v in values) * 1.18
    ax.set_xlim(-xmax, xmax)
    for bar, v in zip(bars, values):
        off = xmax * 0.015
        ax.text(v + (off if v >= 0 else -off), bar.get_y() + bar.get_height() / 2,
                f"{v:+,.0f}", va="center", ha="left" if v >= 0 else "right",
                fontsize=11, color="#333333", fontweight="bold")

    # 标注本策略
    ax.text(xmax * 0.62, -0.34, "本策略 +2243：规避了全交易的 12.8 万大亏", fontsize=10.5,
            color=NAVY_HEX, style="italic")
    ax.spines[["top", "right"]].set_visible(False)
    ax.grid(axis="x", linestyle="--", alpha=0.35)
    fig.savefig(path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print("chart saved:", path)


# ----------------------------------------------------------------------------
# PPT 辅助函数
# ----------------------------------------------------------------------------
def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_rect(slide, left, top, width, height, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, left, top, width, height)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def add_title(slide, text, num=None):
    """标题 + 深蓝左侧色条 + 底部金色细线 + 页脚。"""
    add_rect(slide, MARGIN, Inches(0.38), Inches(0.12), Inches(0.52), NAVY)
    tb = slide.shapes.add_textbox(Inches(0.85), Inches(0.30), Inches(11.6), Inches(0.72))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(27)
    r.font.bold = True
    r.font.color.rgb = NAVY
    # 金色线
    add_rect(slide, MARGIN, Inches(1.12), CONTENT_W, Pt(2.2), GOLD)
    # 页脚
    add_footer(slide, num)


def add_footer(slide, num=None):
    tb = slide.shapes.add_textbox(MARGIN, Inches(7.08), Inches(9), Inches(0.32))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "CA-ISO 电价价差预测 · 讲解版（2026-08-07）"
    r.font.size = Pt(9)
    r.font.color.rgb = GRAY
    if num is not None:
        tb2 = slide.shapes.add_textbox(Inches(12.3), Inches(7.08), Inches(0.9), Inches(0.32))
        p2 = tb2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        r2.text = str(num)
        r2.font.size = Pt(9)
        r2.font.color.rgb = GRAY


def add_bullets(slide, items, left, top, width, height, base_size=17, spacing=8):
    """items: list of (text, level, bold)"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, level, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(spacing)
        p.line_spacing = 1.12
        r = p.add_run()
        r.text = ("•  " if level == 0 else "–  ") + text
        r.font.size = Pt(base_size - 2 * level)
        r.font.bold = bold
        r.font.color.rgb = NAVY if (level == 0 and bold) else DARK
    return tb


def set_cell(cell, text, size=12, bold=False, color=DARK, fill=None,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE):
    cell.margin_left = Inches(0.06)
    cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = valign
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    lines = str(text).split("\n")
    for j, line in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.0
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def add_table(slide, data, left, top, width, col_ratios, row_heights,
              header_fill=NAVY, body_fill=WHITE, banded=True, size=11.5,
              col_align=None, col_bold=None, col_color=None):
    """data: 二维数组（含表头）。col_align/col_bold/col_color: 每列列表或 None。"""
    n_rows, n_cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width,
                                       Inches(sum(row_heights)))
    tbl = tbl_shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    # 列宽
    total = sum(col_ratios)
    for c, ratio in enumerate(col_ratios):
        tbl.columns[c].width = Emu(int(int(width) * ratio / total))
    # 行高
    for r, h in enumerate(row_heights):
        tbl.rows[r].height = Inches(h)

    for r in range(n_rows):
        for c in range(n_cols):
            cell = tbl.cell(r, c)
            is_header = (r == 0)
            if is_header:
                set_cell(cell, data[r][c], size=size, bold=True,
                         color=WHITE, fill=header_fill)
            else:
                fill = body_fill if (not banded or r % 2 == 1) else LIGHT
                al = col_align[c] if col_align else PP_ALIGN.CENTER
                bd = col_bold[c] if col_bold else False
                cc = col_color[c] if col_color else None
                if isinstance(cc, list):
                    cl = cc[r] if cc[r] is not None else DARK
                else:
                    cl = cc if cc is not None else DARK
                set_cell(cell, data[r][c], size=size, bold=bd, color=cl, fill=fill,
                         align=al)
    return tbl_shape


def add_picture_keep_ratio(slide, img_path, left, top, max_w, max_h):
    from PIL import Image
    w, h = Image.open(img_path).size
    ratio = w / h
    if max_w / max_h > ratio:
        width, height = Inches(max_h * ratio), Inches(max_h)
    else:
        width, height = Inches(max_w), Inches(max_w / ratio)
    slide.shapes.add_picture(img_path, left, top, width=width, height=height)
    return width, height


# ----------------------------------------------------------------------------
# 生成图表
# ----------------------------------------------------------------------------
feature_chart_png = os.path.join(CHART_DIR, "feature_importance.png")
baseline_chart_png = os.path.join(CHART_DIR, "baseline_compare.png")
make_feature_chart(feature_chart_png)
make_baseline_chart(baseline_chart_png)

# ----------------------------------------------------------------------------
# 组装 PPT
# ----------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ===== 第 1 页：封面 =====
slide = add_slide(prs)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY_DARK)
add_rect(slide, 0, Inches(5.05), SLIDE_W, Pt(3), GOLD)
# 顶部小字
tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.15), Inches(11.3), Inches(0.5))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "CA-ISO 电价价差预测 · 工程讲解"
r.font.size = Pt(18); r.font.color.rgb = GOLD; r.font.bold = True
# 主标题
tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.4), Inches(1.2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "CA-ISO 电价价差预测"
r.font.size = Pt(52); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "基于 CatBoost 的方向分类与单边套利策略"
r2.font.size = Pt(22); r2.font.color.rgb = RGBColor(0xBF, 0xD3, 0xE6)
# 关键数字
metrics = [
    ("63.9%", "整体方向准确率"),
    ("68.8%", "交易时段命中率 (test)"),
    ("+2243", "test 模拟套利收益"),
    ("109.1", "spread MAE（基线 129.0）"),
]
card_w, card_h = Inches(2.72), Inches(1.35)
gap = Inches(0.24)
start_x = Inches(1.0)
for i, (num, label) in enumerate(metrics):
    x = start_x + i * (card_w + gap)
    card = add_rect(slide, x, Inches(5.5), card_w, card_h, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    card.adjustments[0] = 0.12
    tfb = card.text_frame; tfb.word_wrap = True
    tfb.margin_left = Inches(0.08); tfb.margin_right = Inches(0.08)
    tfb.margin_top = Inches(0.12); tfb.margin_bottom = Inches(0.08)
    pn = tfb.paragraphs[0]; pn.alignment = PP_ALIGN.CENTER
    rn = pn.add_run(); rn.text = num
    rn.font.size = Pt(30); rn.font.bold = True; rn.font.color.rgb = GOLD
    pl = tfb.add_paragraph(); pl.alignment = PP_ALIGN.CENTER
    rl = pl.add_run(); rl.text = label
    rl.font.size = Pt(12); rl.font.color.rgb = WHITE
# 底部日期
tb = slide.shapes.add_textbox(Inches(1.0), Inches(6.85), Inches(11.3), Inches(0.5))
p = tb.text_frame.paragraphs[0]
r = p.add_run(); r.text = "2026-08-07 ｜ 数据对齐 → 特征工程 → 方向分类建模 → 业务评估 → Flask 网页交付"
r.font.size = Pt(13); r.font.color.rgb = RGBColor(0x9F, 0xB8, 0xD0)

# ===== 第 2 页：业务背景 =====
slide = add_slide(prs)
add_title(slide, "业务背景：CA-ISO 价差套利", 2)
add_bullets(slide, [
    ("加州电力市场两阶段定价：日前市场 DA（T-1 撮合）→ 实时市场 RTPD（T 结算），价差 = DA − RTPD", 0, False),
    ("套利动作：T-1 日在日前市场买卖 1 度电，T 日在实时市场反向平仓，赚取价差", 0, False),
    ("三种信号对应的动作：", 0, True),
    ("价差为正（DA 更高）→ 日前卖、实时买", 1, False),
    ("价差为负（RTPD 更高）→ 日前买、实时卖", 1, False),
    ("把握不足 → 观望（hold）", 1, False),
    ("项目目标：预测次日（D+1）各节点、各小时价差方向，据此给出买卖建议", 0, True),
    ("交付：本地 Flask 网页，输入节点 + 日期 → 24h 价差/价格曲线 + 买卖建议 + 回测对比", 0, False),
], Inches(0.85), Inches(1.45), Inches(11.7), Inches(5.4), base_size=17, spacing=10)

# ===== 第 3 页：数据概况 =====
slide = add_slide(prs)
add_title(slide, "数据概况：五个数据源 + 对齐产物", 3)
rows = [
    ["数据", "内容", "时间范围", "说明"],
    ["价格数据/*.xlsx", "3 节点 DA / RTPD / 价差（$/MWh，逐小时）", "2024-01 ~ 2026-08", "目标变量来源"],
    ["load_CA_ISO_TAC_2DA.csv", "官方日前负荷预测（系统级）", "2025-10 ~ 2026-08", "D+1 可直接使用的“未来已知”特征"],
    ["load_CA_ISO_TAC_ACTUAL.csv", "实际负荷（系统级）", "2025-04 ~ 2026-08", "仅作历史滞后特征"],
    ["zone_weather_hourly.csv", "分区天气：温度 / 太阳辐射 / 100m 风速", "2025-04 ~ 2026-08", "视为预报值，D+1 可直接使用"],
    ["节点位置.xlsx", "节点 → 区域映射", "—", "SNLNDRO / CONTROLX → ZP26"],
]
add_table(slide, rows, Inches(0.55), Inches(1.45), Inches(12.2),
          [2.1, 4.2, 2.0, 3.9], [0.42, 0.62, 0.62, 0.62, 0.62, 0.5], size=11.5)
add_bullets(slide, [
    ("建模节点：SNLNDRO、CONTROLX（ZP26，已建模）", 0, True),
    ("ELCAJNGT（SP15）：数据仅 2026-03 起、训练窗口 0 样本，未建模", 0, False),
    ("对齐产物：master.csv（45780 行长表）→ features.parquet（特征矩阵）", 0, False),
], Inches(0.85), Inches(5.35), Inches(11.7), Inches(1.5), base_size=15, spacing=6)

# ===== 第 4 页：预测方法（1）任务定义 + 防泄漏 =====
slide = add_slide(prs)
add_title(slide, "预测方法（1）：任务定义与防泄漏", 4)
add_bullets(slide, [
    ("任务：每个样本 = 决策日 D 的某个小时，目标 = 预测 D+1 同小时的价差方向", 0, False),
    ("防泄漏红线：决策时点 D 日 10:00 PT 前，D 日当天的实际价格/负荷不可见", 0, True),
    ("滞后特征只用 D-1 及更早；D+1 的 2DA 负荷与天气是预报值、此时已可得，可作特征", 1, False),
    ("时间切分（按决策日，不随机）：", 0, True),
    ("train 2025-04~12（13,700 样本）", 1, False),
    ("val 2026-01~05（7,550 样本）", 1, False),
    ("test 2026-06~08（3,300 样本）", 1, False),
    ("特征工程：共 34 个特征，分五类 —— 历史滞后 / 滚动统计 / 未来已知 / 日历 / 节点", 0, False),
], Inches(0.85), Inches(1.5), Inches(11.7), Inches(5.3), base_size=16.5, spacing=9)

# ===== 第 5 页：预测方法（2）模型 + 决策规则 =====
slide = add_slide(prs)
add_title(slide, "预测方法（2）：方向分类器 + 单边决策规则", 5)
add_bullets(slide, [
    ("决策核心：CatBoost 方向分类器，目标 y = (价差 > 0)，直接输出 P(价差 > 0)", 0, True),
    ("训练用 val 早停，缺省特征由 CatBoost 原生处理 NaN", 0, False),
    ("展示辅助：分位数回归（q10 / q50 / q90）画价差曲线与置信带，DA/RTPD q50 画价格曲线", 0, False),
    ("实测 CatBoost 优于 LightGBM：test 套利收益 +30%（+1744 → +2243）", 0, True),
    ("单边决策（无买）：sell 约 33% / hold 约 67% —— 规则经回测校准：", 0, True),
], Inches(0.85), Inches(1.5), Inches(11.7), Inches(3.4), base_size=16.5, spacing=9)
# 规则框
rule_box = add_rect(slide, Inches(1.35), Inches(4.55), Inches(10.6), Inches(1.35),
                    LIGHT_BLUE, line=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rule_box.adjustments[0] = 0.08
tf = rule_box.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.16)
p = tf.paragraphs[0]; p.line_spacing = 1.3
r = p.add_run()
r.text = "P(价差 > 0) > 0.5   且   近 7 日价差波动 std7 ≤ 120   →   “卖”（日前卖、实时买）"
r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = NAVY
p2 = tf.add_paragraph(); p2.line_spacing = 1.2
r2 = p2.add_run(); r2.text = "否则  →  “观望”（hold）"
r2.font.size = Pt(15); r2.font.color.rgb = DARK

# ===== 第 6 页：各节点结果 =====
slide = add_slide(prs)
add_title(slide, "各节点结果：SNLNDRO 是策略主力", 6)
rows = [
    ["节点", "时段", "交易数(占比)", "交易命中率", "单边策略收益", "实际正价差占比", "价差幅度 |实际|"],
    ["SNLNDRO", "test", "1018 (65%)", "68.8%", "+2243", "62.6%", "7.8"],
    ["CONTROLX", "test", "0 (0%)",    "—",    "0",     "33.9%", "213.8"],
    ["整体",     "val",  "—",          "71.3%", "+6884", "—",     "—"],
    ["",         "test", "—",          "68.8%", "+2243", "—",     "—"],
]
col_color = [None, None, None, None,
             [None, GREEN, DARK, GREEN, GREEN],
             None, None]
add_table(slide, rows, Inches(0.55), Inches(1.42), Inches(12.2),
          [1.5, 0.9, 1.7, 1.6, 1.7, 1.9, 1.5], [0.42, 0.52, 0.52, 0.52, 0.52],
          size=12, col_align=[PP_ALIGN.CENTER]*7)
# 合并节点列（整体 val/test）
tbl = slide.shapes[-1].table
tbl.cell(3, 0).merge(tbl.cell(4, 0))
set_cell(tbl.cell(3, 0), "整体", size=13, bold=True, color=NAVY,
         fill=WHITE if 3 % 2 else LIGHT)
add_bullets(slide, [
    ("SNLNDRO：价差系统性偏正（test 实际正价差 62.6%），预测正价差命中 68.8%，幅度小（±8）风险可控 → 稳定盈利", 0, True),
    ("CONTROLX：实际正价差仅占 33.9%、幅度高达 ±214，CatBoost 几乎不给出“卖”信号（test 0 笔交易）→ 主动避开大幅度节点", 0, False),
    ("ELCAJNGT 未建模（数据不足）", 0, False),
], Inches(0.85), Inches(5.55), Inches(11.7), Inches(1.4), base_size=14, spacing=5)

# ===== 第 7 页：影响因子 =====
slide = add_slide(prs)
add_title(slide, "影响因子：DA 价滞后 + 价差波动主导", 7)
w, h = add_picture_keep_ratio(slide, feature_chart_png, Inches(0.9), Inches(1.5),
                              Inches(11.6), Inches(4.6))
add_bullets(slide, [
    ("历史滞后共 44.4%（da_lag2 15.8% + da_lag1 5.5% 领衔）—— 日前撮合价对次日方向最具预示力", 0, True),
    ("波动统计升至第二（24.1%：spread_std14 13.1%、spread_std7 7.8%）；未来预报 22.9%、日历 9.7%、节点 6.9%", 0, False),
], Inches(0.85), Inches(6.2), Inches(11.7), Inches(0.85), base_size=14, spacing=4)

# ===== 第 8 页：基线对比 =====
slide = add_slide(prs)
add_title(slide, "基线对比：单边策略显著优于全交易", 8)
w, h = add_picture_keep_ratio(slide, baseline_chart_png, Inches(0.9), Inches(1.5),
                              Inches(11.6), Inches(4.2))
add_bullets(slide, [
    ("单边策略 +2243 vs 全交易 -128,259：规避了约 12.8 万美元的大亏", 0, True),
    ("反向 +128,259 为“全交易取负”的镜像上界，非可执行策略；不交易 0 为下限", 0, False),
    ("数值基线（spread 预测）：naive（D-1 同小时直推）MAE = 129.0 → 本模型 MAE = 109.1（提升 15%）", 0, True),
], Inches(0.85), Inches(6.0), Inches(11.7), Inches(1.1), base_size=14.5, spacing=4)

# ===== 第 9 页：套利收益曲线 =====
slide = add_slide(prs)
add_title(slide, "套利收益曲线：test 累计 +2243", 9)
add_picture_keep_ratio(slide, ARB_PNG, Inches(1.35), Inches(1.55), Inches(10.6), Inches(4.35))
add_bullets(slide, [
    ("test（2026-06 ~ 08，65 天）模拟套利累计 +2243，日均约 +35", 0, True),
    ("最大单笔亏损 -171（单边规则 + std7 风控下亏损可控）", 0, False),
    ("val 同时段为正：+6884（2026 H1）—— 双时段验证，避免只在测试集上调参", 0, False),
], Inches(0.85), Inches(6.1), Inches(11.7), Inches(1.0), base_size=14.5, spacing=4)

# ===== 第 10 页：评测方法 =====
slide = add_slide(prs)
add_title(slide, "评测方法：四层指标，业务优先", 10)
add_bullets(slide, [
    ("① 方向准确率（核心业务指标）：整体 63.9%（val 65.3%）；交易时段命中率 68.8%（val 71.3%）（只看实际下单样本，更贴业务）", 0, True),
    ("② 数值精度（参考）：spread MAE 109.1 vs naive 129.0；DA/RTPD MAE、RMSE —— 价差方向才是业务目标", 0, False),
    ("③ 决策质量：卖 33% / 观望 67%（单边、无买）；每类决策平均收益；hold 是否避开了大额亏损", 0, False),
    ("④ 模拟套利（业务落地）：按预测方向 DA 买卖、RTPD 反向平仓、按真实价差结算 → 总收益、日均、命中率、最大单笔亏损、累计曲线", 0, True),
    ("⑤ 稳健性：val（2026 H1）+6884 与 test（2026 H2）+2243 双时段均正收益", 0, True),
], Inches(0.85), Inches(1.6), Inches(11.7), Inches(5.2), base_size=16.5, spacing=11)

# ===== 第 11 页：增强方向 =====
slide = add_slide(prs)
add_title(slide, "增强方向：需外部数据（按优先级）", 11)
rows = [
    ["外部数据", "为什么有用", "当前缺失"],
    ["天然气价格（Henry Hub / 加州气价）", "燃气机组是加州边际定价者，电价 ≈ 气价 × 热耗率；DA/RTPD 用不同时点气价 → 直接是价差方向的供给侧信号", "历史价滞后只间接含过去气价，无 D+1 气价预测"],
    ["CA-ISO 光伏/风电出力预测", "可再生能源出力直接影响 RTPD 与价差（鸭子曲线核心驱动）", "ssrd 只是代理，非实际出力"],
    ["负荷偏差（实时 vs 日前预测）", "日内偏差是 RTPD 波动的直接诱因", "缺"],
    ["更多节点/区域价格", "阻塞与区域价差预示方向", "仅 2 节点"],
    ["滚动 / 更长历史训练", "适应市场漂移", "已扩至 9 个月，仍可加滚动窗口"],
]
add_table(slide, rows, Inches(0.55), Inches(1.45), Inches(12.2),
          [2.7, 5.7, 3.8], [0.42, 0.82, 0.82, 0.62, 0.62, 0.62], size=11.5,
          col_align=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])

# ===== 第 12 页：已知局限 =====
slide = add_slide(prs)
add_title(slide, "已知局限（如实说明）", 12)
add_bullets(slide, [
    ("套利利润偏薄：val +6884 / test +2243（test 日均约 +35）—— 价差本身难预测，63.9% 方向准确率带来的收益空间有限", 0, True),
    ("单边牺牲双边机会：只做“卖”（正价差）方向，负价差方向的套利机会被放弃", 0, False),
    ("交易完全集中在 SNLNDRO：CONTROLX 在单边规则下 test 0 笔交易", 0, False),
    ("跨时段仍有漂移：val / test 收益均为正但绝对值不同，市场状态变化仍影响表现", 0, False),
    ("决策阈值基于回测校准（prob > 0.5、std7 ≤ 120），实盘需再验证", 0, False),
    ("ELCAJNGT 未训练；天气按“预报值”假设（若用实测需改滞后特征）", 0, False),
], Inches(0.85), Inches(1.6), Inches(11.7), Inches(5.2), base_size=16.5, spacing=11)

# ===== 第 13 页：结语 =====
slide = add_slide(prs)
add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY_DARK)
add_rect(slide, 0, Inches(4.15), SLIDE_W, Pt(3), GOLD)
tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.1))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run(); r.text = "小结"
r.font.size = Pt(18); r.font.color.rgb = GOLD; r.font.bold = True
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "已交付一套可运行、可验证的价差方向预测与单边套利系统"
r2.font.size = Pt(36); r2.font.bold = True; r2.font.color.rgb = WHITE

add_bullets(slide, [
    ("数据 → 特征 → 模型 → 业务评估 → 网页交付，全流程闭环；防泄漏严格、按决策日切分", 0, True),
    ("方向准确率 63.9%（交易时段 68.8%），模拟套利 val +6884 / test +2243，收益为正、风险可控", 0, True),
    ("下一步：接入气价 / 新能源出力预测等外部数据，扩展节点与滚动训练", 0, False),
], Inches(1.0), Inches(3.0), Inches(11.3), Inches(1.9), base_size=17, spacing=9)

# 联系方式占位
card = add_rect(slide, Inches(1.0), Inches(5.15), Inches(11.3), Inches(1.5),
                NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
card.adjustments[0] = 0.1
tf = card.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "联系人：＿＿＿＿＿＿＿　｜　邮箱：＿＿＿＿＿@＿＿＿"
r.font.size = Pt(16); r.font.color.rgb = WHITE
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "谢谢观看 · 欢迎提问"
r2.font.size = Pt(14); r2.font.color.rgb = GOLD

prs.save(OUT_PATH)
print("PPT saved:", OUT_PATH)
print("size bytes:", os.path.getsize(OUT_PATH))
